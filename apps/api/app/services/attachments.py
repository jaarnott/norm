"""Chat attachments — turning an uploaded file into something the model reads.

A user can attach files to a chat message. PDFs and images go to the model as
native Anthropic ``document``/``image`` blocks (base64); Office and text files
are extracted to text server-side (the model can't read a .docx/.xlsx byte
stream). The bytes live in ``UploadedDocument``; this module is the one place
that turns those bytes into a content block, and lists a thread's attachments so
the model can re-open one with the ``get_attachment`` tool.
"""

from __future__ import annotations

import base64
import io
import logging

logger = logging.getLogger(__name__)

# Media types Claude reads directly — sent as a native block, no extraction.
_NATIVE_TYPES = frozenset(
    {
        "application/pdf",
        "image/png",
        "image/jpeg",
        "image/jpg",
        "image/gif",
        "image/webp",
    }
)

_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

#: What the chat composer may attach — native block types plus everything we can
#: extract to text. Used to gate the upload/attach path.
CHAT_ALLOWED_TYPES = _NATIVE_TYPES | {
    _DOCX,
    _XLSX,
    _PPTX,
    "text/plain",
    "text/csv",
    "text/markdown",
    "application/json",
}

_TEXT_EXTS = {
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".log",
    ".yaml",
    ".yml",
    ".py",
    ".ts",
    ".tsx",
    ".js",
    ".jsx",
    ".sql",
    ".html",
    ".css",
}

#: A single extracted file is capped so one huge spreadsheet can't blow out the
#: context window.
_MAX_TEXT_CHARS = 200_000

#: Anthropic rejects any single image over ~5MB and treats ~1568px as the largest
#: useful long edge. A raw phone photo (prod thread eb83d8f0, 24 Aug 2026: a 7.78MB
#: PXL_*.jpg) blows past both and fails the whole request. We normalise images to
#: sit inside these limits before sending — the same thing claude.ai does.
_IMAGE_MAX_EDGE = 1568
_IMAGE_MAX_BYTES = 4_500_000  # safety margin under Anthropic's ~5MB per-image cap


def _ext(filename: str | None) -> str:
    if filename and "." in filename:
        return "." + filename.rsplit(".", 1)[-1].lower()
    return ""


def _kind(content_type: str | None, filename: str | None) -> tuple[str, str]:
    """Classify a file into how we turn it into a block. Content-type is trusted
    first (browsers set it), extension is the fallback when it's generic."""
    ct = (content_type or "").split(";")[0].strip().lower()
    ext = _ext(filename)
    if ct in _NATIVE_TYPES:
        return "native", ct
    if ct == _DOCX or ext == ".docx":
        return "docx", ct
    if ct == _XLSX or ext == ".xlsx":
        return "xlsx", ct
    if ct == _PPTX or ext == ".pptx":
        return "pptx", ct
    if ct.startswith("image/"):
        return "native", ct
    if ct == "application/pdf" or ext == ".pdf":
        return "native", "application/pdf"
    if ct.startswith("text/") or ct == "application/json" or ext in _TEXT_EXTS:
        return "text", ct or "text/plain"
    return "unknown", ct


def _truncate(text: str) -> str:
    if len(text) > _MAX_TEXT_CHARS:
        return (
            text[:_MAX_TEXT_CHARS]
            + "\n\n[… truncated — the file is larger than the extraction limit …]"
        )
    return text


def _extract_docx(data: bytes) -> str:
    import docx

    d = docx.Document(io.BytesIO(data))
    parts = [p.text for p in d.paragraphs if p.text]
    for table in d.tables:
        for row in table.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out: list[str] = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                out.append("\t".join("" if c is None else str(c) for c in row))
    wb.close()
    return "\n".join(out)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs)
                    if txt:
                        out.append(txt)
    return "\n".join(out)


def _encode_jpeg(img) -> bytes:
    """Flatten to RGB and JPEG-encode. Starts at high quality (q90 — the size after
    downscaling to 1568px is tiny, so fidelity is nearly free) and only steps down if
    an unusually busy image would exceed the byte cap."""
    rgb = img.convert("RGB")
    buf = io.BytesIO()
    for quality in (90, 80, 70, 55):
        buf.seek(0)
        buf.truncate(0)
        rgb.save(buf, format="JPEG", quality=quality, optimize=True)
        if buf.tell() <= _IMAGE_MAX_BYTES:
            break
    return buf.getvalue()


def normalize_image_for_anthropic(
    data: bytes, content_type: str | None
) -> tuple[bytes, str]:
    """Downscale / re-encode an image so it fits Anthropic's per-image limits.

    Anthropic rejects a single image over ~5MB and treats ~1568px as the largest
    useful long edge; a raw phone photo blows past both and fails the whole Messages
    request. This mirrors what claude.ai does for you: fix EXIF orientation, downscale
    the long edge to 1568px, and re-encode so the bytes land under the cap. Small,
    in-bounds, correctly-oriented images are returned untouched.

    Returns ``(bytes, media_type)``. Never raises — on any failure it logs and returns
    the original bytes, so behaviour is no worse than before.
    """
    ct = (content_type or "").split(";")[0].strip().lower()
    try:
        from PIL import Image, ImageOps

        img = Image.open(io.BytesIO(data))
        img.load()
        # Phones bake rotation into EXIF orientation (tag 0x0112), which Anthropic
        # ignores. Detect it from the tag (a value other than 1) — reading it here
        # is reliable across Pillow versions, unlike checking exif_transpose's return.
        try:
            orientation = img.getexif().get(0x0112, 1)
        except Exception:
            orientation = 1
        reoriented = orientation not in (0, 1)
        img = ImageOps.exif_transpose(img) or img  # no-op when orientation is 1

        oversized = max(img.size) > _IMAGE_MAX_EDGE or len(data) > _IMAGE_MAX_BYTES
        if not oversized and not reoriented:
            return data, ct or "image/jpeg"

        if max(img.size) > _IMAGE_MAX_EDGE:
            # LANCZOS is the highest-quality downscaling filter — worth it since we
            # only do this once per image and detail matters (comparing two dishes).
            img.thumbnail((_IMAGE_MAX_EDGE, _IMAGE_MAX_EDGE), Image.LANCZOS)

        has_alpha = img.mode in ("RGBA", "LA") or (
            img.mode == "P" and "transparency" in img.info
        )
        if has_alpha:
            buf = io.BytesIO()
            img.convert("RGBA").save(buf, format="PNG", optimize=True)
            if buf.tell() <= _IMAGE_MAX_BYTES:
                return buf.getvalue(), "image/png"
            # Transparent but still too big — flatten to JPEG (loses alpha).
        return _encode_jpeg(img), "image/jpeg"
    except Exception:
        logger.warning(
            "image normalise failed (type=%s) — sending original", ct, exc_info=True
        )
        return data, ct or "application/octet-stream"


def build_content_block(
    data: bytes, content_type: str | None, filename: str | None
) -> dict:
    """One Anthropic content block for an attachment.

    PDF/image → a native base64 ``document``/``image`` block. docx/xlsx/pptx and
    text files → an extracted-text block labelled with the filename. Raises
    ``ValueError`` for a type we can neither send nor extract.
    """
    kind, ct = _kind(content_type, filename)
    name = filename or "file"
    if kind == "native":
        if ct.startswith("image/"):
            # Fit the image inside Anthropic's per-image limits (a large phone photo
            # would otherwise fail the whole request).
            data, ct = normalize_image_for_anthropic(data, ct)
        b64 = base64.b64encode(data).decode()
        return {
            "type": "image" if ct.startswith("image/") else "document",
            "source": {"type": "base64", "media_type": ct, "data": b64},
        }
    if kind == "text":
        text = data.decode("utf-8", errors="replace")
    elif kind == "docx":
        text = _extract_docx(data)
    elif kind == "xlsx":
        text = _extract_xlsx(data)
    elif kind == "pptx":
        text = _extract_pptx(data)
    else:
        raise ValueError(f"Can't read {ct or 'this file type'}.")
    return {"type": "text", "text": f"[Attachment: {name}]\n{_truncate(text)}"}


def link_chat_attachments(
    attachment_ids: list[str], thread_id: str, user_id: str | None, db
) -> list[dict]:
    """Stamp uploaded files with the thread they were attached to and return
    their reference metadata for the user Message. Only the uploader's own files
    are linked — an id that belongs to someone else (or doesn't exist) is
    dropped. Order is preserved."""
    if not attachment_ids:
        return []
    from app.db.models import UploadedDocument

    by_id = {
        d.id: d
        for d in db.query(UploadedDocument)
        .filter(UploadedDocument.id.in_(attachment_ids))
        .all()
    }
    out: list[dict] = []
    for aid in attachment_ids:
        d = by_id.get(aid)
        if not d:
            continue
        if user_id and d.user_id and d.user_id != user_id:
            continue
        d.thread_id = thread_id
        out.append(
            {
                "upload_id": d.id,
                "filename": d.filename,
                "content_type": d.content_type,
                "size": d.size,
            }
        )
    return out


def attachment_manifest(thread, db) -> str | None:
    """A compact list of every file attached in this thread, so the model knows
    what it can re-open with ``get_attachment`` on a later turn."""
    if thread is None or db is None:
        return None
    from app.db.models import Message

    msgs = (
        db.query(Message)
        .filter(Message.thread_id == thread.id, Message.attachments.isnot(None))
        .all()
    )
    items: list[str] = []
    seen: set[str] = set()
    for m in msgs:
        for a in m.attachments or []:
            uid = a.get("upload_id")
            if uid and uid not in seen:
                seen.add(uid)
                items.append(f"- {a.get('filename') or 'file'} (id: {uid})")
    if not items:
        return None
    return (
        "[Attachments in this conversation]\n"
        "Files the user attached — call get_attachment with an id to re-open one:\n"
        + "\n".join(items)
    )
